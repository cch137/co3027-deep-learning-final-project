"""Generate the final-project deck: outputs/LeadLagNet_Final.pptx (16:9, English).

Run:  python gen_pptx.py
Placeholders the user must fill are marked with [[ ... ]].
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

FIG = Path("outputs/figures")
OUT = Path("outputs/LeadLagNet_Final.pptx")

NAVY = RGBColor.from_string("1B2A4A")
INK = RGBColor.from_string("22304A")
AMBER = RGBColor.from_string("F59E0B")
LIGHT = RGBColor.from_string("F4F6FA")
WHITE = RGBColor.from_string("FFFFFF")
GRAY = RGBColor.from_string("64748B")
LINE = RGBColor.from_string("D9DFEA")
GREEN = RGBColor.from_string("0B7A55")
RED = RGBColor.from_string("B91C1C")

W, H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
slide_no = 0


def new_slide(dark: bool = False):
    global slide_no
    slide_no += 1
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY if dark else LIGHT
    bg.line.fill.background()
    bg.shadow.inherit = False
    # motif: left navy band + amber slide-number chip
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.30), H)
    band.fill.solid()
    band.fill.fore_color.rgb = AMBER if dark else NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    chip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.25), Inches(0.30), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = AMBER
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(slide_no)
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = FONT
    return s


def textbox(s, x, y, w, h, lines, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
            leading=1.12, anchor=MSO_ANCHOR.TOP):
    """lines: list of str or (str, dict-overrides)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in lines:
        text, ov = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("leading", leading)
        p.space_after = Pt(ov.get("space_after", 4))
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
    return box


def kicker_title(s, kicker, title, title_size=27):
    textbox(s, Inches(0.62), Inches(0.22), Inches(12.2), Inches(0.32),
            [(kicker.upper(), {"size": 12, "bold": True, "color": AMBER})])
    textbox(s, Inches(0.62), Inches(0.50), Inches(12.3), Inches(0.75),
            [(title, {"size": title_size, "bold": True, "color": NAVY})])


def card(s, x, y, w, h, fill=WHITE, line_color=LINE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = 0.045
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = line_color
    c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


def picture(s, path, x, y, w, h, frame=True):
    """Fit image inside (x,y,w,h) preserving aspect, centered, optional card behind."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = int(x + (w - pw) / 2), int(y + (h - ph) / 2)
    if frame:
        card(s, Emu(px - 60000), Emu(py - 60000), Emu(pw + 120000), Emu(ph + 120000))
    s.shapes.add_picture(str(path), Emu(px), Emu(py), Emu(pw), Emu(ph))


def flow_box(s, x, y, w, h, title, sub, fill=NAVY, title_color=WHITE, sub_color=None):
    b = card(s, x, y, w, h, fill=fill, line_color=fill)
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(12.5)
    r.font.bold = True
    r.font.color.rgb = title_color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT
        r2.font.size = Pt(10)
        r2.font.color.rgb = sub_color or RGBColor.from_string("C7D2E5")
    return b


def arrow(s, x, y, w=Inches(0.28)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.22))
    a.fill.solid()
    a.fill.fore_color.rgb = AMBER
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def bullets(items, size=13.5):
    out = []
    for it in items:
        if isinstance(it, tuple):
            text, ov = it
            ov = {"size": size, **ov}
        else:
            text, ov = it, {"size": size}
        out.append(("•  " + text, ov))
    return out


# ============================================================ 1. title
s = new_slide(dark=True)
textbox(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.4),
        [("FINAL PROJECT REPORT", {"size": 13, "bold": True, "color": AMBER})])
textbox(s, Inches(0.9), Inches(1.85), Inches(11.8), Inches(1.55),
        [("Cross-Market Lead-Lag Estimation", {"size": 32, "bold": True, "color": WHITE}),
         ("for Financial Assets Using Deep Learning", {"size": 32, "bold": True, "color": WHITE})])
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.62), Inches(1.6), Inches(0.045))
ln.fill.solid(); ln.fill.fore_color.rgb = AMBER; ln.line.fill.background(); ln.shadow.inherit = False
textbox(s, Inches(0.9), Inches(3.95), Inches(11.4), Inches(0.9),
        [("A neural estimator of lag, correlation strength, and scale between financial "
          "time series, trained with synthetic supervision.",
          {"size": 17, "color": RGBColor.from_string("D7DEED")})])
textbox(s, Inches(0.9), Inches(5.15), Inches(11.4), Inches(1.5),
        [("徐崇恆  112503023", {"size": 16, "bold": True, "color": WHITE}),
         ("Programming for Deep Learning   •   2026/06/12",
          {"size": 13, "color": RGBColor.from_string("9FB0CC")})])

# ============================================================ 2. background
s = new_slide()
kicker_title(s, "Background", "Lead-Lag Relationships in Time Series")
textbox(s, Inches(0.62), Inches(1.45), Inches(5.1), Inches(5.6), bullets([
    ("One market often moves first and another follows: a lead-lag relationship.", {}),
    ("Classical tool — the cross-correlation scan:", {}),
    ("    ρ(τ) = corr( A(t−τ), B(t) ),  pick τ* = argmax |ρ(τ)|", {"size": 13, "color": GRAY}),
    ("It assumes the link is linear and stationary, and it must re-scan every lag, "
     "every pair, every statistic you care about.", {}),
    ("Real example (right): the S&P 500 close leads Asian indices by exactly one "
     "trading day — a structural, timezone-driven effect (corr 0.42 forward vs 0.004 reverse).", {}),
    ("Question of this project: can a small neural network learn the skill of "
     "reading this relationship directly from two raw windows?", {"bold": True}),
]))
picture(s, FIG / "fig3_leadlag_motivation.png", Inches(5.95), Inches(1.5), Inches(7.1), Inches(5.5))

# ============================================================ 3. problem
s = new_slide()
kicker_title(s, "Problem", "Problem Statement & Model I/O")
c = card(s, Inches(0.62), Inches(1.5), Inches(12.1), Inches(1.25), fill=NAVY, line_color=NAVY)
tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
p = tf.paragraphs[0]; r = p.add_run()
r.text = ("Given two aligned time series and zero human annotations, estimate who leads whom, "
          "by how many steps, and how strongly — including sign and scale.")
r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE
card(s, Inches(0.62), Inches(3.05), Inches(5.9), Inches(3.6))
textbox(s, Inches(0.92), Inches(3.25), Inches(5.3), Inches(3.3),
        [("INPUT", {"size": 12, "bold": True, "color": AMBER})] + bullets([
            ("Two windows, shape (2, 128) — same timestamps, same step size", {}),
            ("Stationary increments (daily log returns; Δ for yields)", {}),
            ("Per-window z-scored — the model sees shape, not scale", {}),
            ("No constraint on era or frequency", {}),
        ], size=13))
card(s, Inches(6.82), Inches(3.05), Inches(5.9), Inches(3.6))
textbox(s, Inches(7.12), Inches(3.25), Inches(5.3), Inches(3.3),
        [("OUTPUT", {"size": 12, "bold": True, "color": AMBER})] + bullets([
            ("Lag: 21-class distribution over τ ∈ [−10, +10]  (τ>0 ⇒ A leads B)", {}),
            ("Correlation ρ̂ ∈ [−1, 1]: sign = direction, |ρ̂|<0.2 ⇒ unrelated", {}),
            ("Scale β̂ = ρ̂·σB/σA — recovered analytically, no extra head", {}),
            ("Lag as classification, not regression: posteriors are multimodal "
             "and the softmax provides uncertainty for free", {}),
        ], size=13))

# ============================================================ 4. key idea
s = new_slide()
kicker_title(s, "Key idea", "Synthetic Supervision: Training Without Manual Annotation")
y = Inches(1.75)
flow_box(s, Inches(0.7), y, Inches(2.65), Inches(1.15), "GENERATOR",
         "B = β·shift(A, τ) + σ·ε — τ, ρ, β known by construction")
arrow(s, Inches(3.43), Inches(2.21))
flow_box(s, Inches(3.80), y, Inches(2.65), Inches(1.15), "~1M LABELED PAIRS",
         "no human annotation, no real-market labels")
arrow(s, Inches(6.53), Inches(2.21))
flow_box(s, Inches(6.90), y, Inches(2.65), Inches(1.15), "TRAIN LeadLagNet",
         "supervised cross-entropy + MSE objectives")
arrow(s, Inches(9.63), Inches(2.21))
flow_box(s, Inches(10.00), y, Inches(2.65), Inches(1.15), "TEST ON REAL MARKETS",
         "no real series used in training", fill=AMBER, title_color=NAVY,
         sub_color=RGBColor.from_string("5B4A1A"))
textbox(s, Inches(0.7), Inches(3.25), Inches(12.0), Inches(0.9), bullets([
    ("Training is fully supervised, but all supervision derives from the generative process — "
     "no manual annotation is involved at any stage.", {}),
    ("At inference the model receives only the two raw windows: no statistics, no lag scanning, "
     "a single forward pass (~2 ms).", {}),
], size=13.5))
principles = [
    ("ZERO ANNOTATION COST", "Every label (τ, ρ, β) is a byproduct of how the pair was "
     "constructed — no human ever labels anything."),
    ("UNLIMITED FRESH DATA", "~1M unique pairs generated on the fly; none stored, so the "
     "model can never overfit a fixed dataset."),
    ("CONTROLLABLE DIFFICULTY", "Noise level σ dials correlation from near-perfect to "
     "invisible — the curriculum covers the whole identifiability range."),
]
x = Inches(0.7)
for title, body in principles:
    card(s, x, Inches(4.45), Inches(3.85), Inches(2.45))
    textbox(s, x + Inches(0.25), Inches(4.68), Inches(3.35), Inches(2.0),
            [(title, {"size": 12, "bold": True, "color": AMBER, "space_after": 8}),
             (body, {"size": 12.5})])
    x += Inches(4.08)

# ============================================================ 5. pipeline
s = new_slide()
kicker_title(s, "Engineering", "System Pipeline")
y = Inches(2.0)
stages = [
    ("DATA DOWNLOAD", "yfinance, 14 assets × 10 yr, one call — no scraping"),
    ("VALIDATOR", "86 automated checks: gaps, freshness, alignment, sanity"),
    ("SYNTH GENERATOR", "~1M pairs on the fly, zero storage"),
    ("LeadLagNet", "130,502-parameter estimator"),
    ("EVAL SUITE", "synthetic / real / out-of-distribution"),
    ("FastAPI DEMO", "bilingual web UI, 2 ms inference"),
]
bw, gap = Inches(1.86), Inches(0.32)
x = Inches(0.62)
for i, (t, sub) in enumerate(stages):
    flow_box(s, x, y, bw, Inches(1.5), t, sub, fill=NAVY if i < 5 else AMBER,
             title_color=WHITE if i < 5 else NAVY,
             sub_color=None if i < 5 else RGBColor.from_string("5B4A1A"))
    if i < 5:
        arrow(s, x + bw + Inches(0.03), y + Inches(0.64), w=Inches(0.26))
    x += bw + gap
textbox(s, Inches(0.62), Inches(4.0), Inches(12.1), Inches(2.9), bullets([
    ("Every stage is a standalone script driven by one shared config.py (universe, window, lags, split dates).", {}),
    ("Strict time discipline end-to-end: splits are applied to the window END date — train ≤ 2025-12-31, "
     "validation = 2026 Q1, test = 2026 Q2; no shuffling across time, normalization uses no future data.", {}),
    ("Fully reproducible: fixed seeds, cached datasets, free data sources, configuration-driven stages.", {}),
    ("The classical cross-correlation baseline was implemented before any model work — it defines the "
     "reference answers all later stages are judged against.", {}),
], size=14))

# ============================================================ 6. data engineering
s = new_slide()
kicker_title(s, "Engineering", "Data: 14 Markets, 10 Years, Leak-Free Splits")
textbox(s, Inches(0.62), Inches(1.45), Inches(5.0), Inches(5.7), bullets([
    ("US equities, volatility, Asian & European indices, crypto, commodities, rates, USD index "
     "(2016 → 2026-06, daily).", {}),
    ("Pitfall found & fixed: crypto trades 7 days/week — a naive union calendar pads equity series "
     "with ~29% fake zero-return weekend rows, distorting every correlation. Calendar restricted to "
     "weekdays; crypto weekend moves fold into Monday.", {}),
    ("Holidays: forward-fill capped at 3 days; longer gaps stay NaN and the window is skipped.", {}),
    ("Known quirks handled explicitly: ETH listing (2017-11), negative oil prices (2020-04).", {}),
    ("86 automated validation checks pass, including end-to-end alignment tests against known market "
     "facts before any model exists.", {"bold": True}),
]))
picture(s, FIG / "fig2_dataset_overview.png", Inches(5.85), Inches(1.5), Inches(7.2), Inches(5.5))

# ============================================================ 7. generator design
s = new_slide()
kicker_title(s, "Method", "Synthetic Generator — the Actual Training Set")
c = card(s, Inches(0.62), Inches(1.5), Inches(12.1), Inches(0.95), fill=NAVY, line_color=NAVY)
tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "B(t) = β · A(t − τ) + σ · ε(t)        labels:  τ (lag),  ρ = realized Pearson at true τ,  β (scale)"
r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
textbox(s, Inches(0.62), Inches(2.7), Inches(5.6), Inches(4.3), bullets([
    ("Three base processes so the model cannot memorize one shape family: white noise, AR(1) "
     "with random φ, mixtures of 2–3 incommensurate sines.", {}),
    ("τ uniform over [−10, +10]; |β| log-uniform in [0.3, 3], both signs; σ log-uniform in [0.1, 3] "
     "— covers near-perfect to nearly invisible correlation.", {}),
    ("15% of pairs are fully independent: teaches the model to answer “unrelated”; their lag loss is masked.", {}),
    ("ρ label = realized correlation of this noise draw, not the theoretical value.", {}),
    ("Cross-checked before training: the classical scan recovers generator labels with 94% lag accuracy "
     "— labels and conventions are consistent.", {"bold": True}),
], size=13.5))
picture(s, FIG / "fig4_synthetic_examples.png", Inches(6.45), Inches(2.75), Inches(6.6), Inches(4.3))

# ============================================================ 8. model design
s = new_slide()
kicker_title(s, "Method", "LeadLagNet — Siamese CNN with Correlation Fusion")
x0, bw2 = Inches(0.72), Inches(4.5)
flow_box(s, x0, Inches(1.5), Inches(2.15), Inches(0.7), "series A (128)", None, fill=WHITE, title_color=NAVY)
flow_box(s, x0 + Inches(2.35), Inches(1.5), Inches(2.15), Inches(0.7), "series B (128)", None, fill=WHITE, title_color=NAVY)
flow_box(s, x0, Inches(2.4), Inches(2.15), Inches(0.85), "shared encoder",
         "Conv1d 1→24 k7 · Conv1d 24→48 k5 s2")
flow_box(s, x0 + Inches(2.35), Inches(2.4), Inches(2.15), Inches(0.85), "shared encoder",
         "same weights (siamese)")
flow_box(s, x0, Inches(3.45), bw2, Inches(0.85), "correlation fusion",
         "concat [ F_A , F_B , F_A ⊙ F_B ] — 144 channels")
flow_box(s, x0, Inches(4.5), bw2, Inches(0.85), "fusion convs + pooling",
         "Conv 144→96 s2 · Conv 96→96 s2 · mean+max pool · FC 192→128")
flow_box(s, x0, Inches(5.55), Inches(2.15), Inches(0.9), "lag head",
         "21-way softmax over τ", fill=AMBER, title_color=NAVY,
         sub_color=RGBColor.from_string("5B4A1A"))
flow_box(s, x0 + Inches(2.35), Inches(5.55), Inches(2.15), Inches(0.9), "ρ head",
         "tanh ∈ [−1, 1]", fill=AMBER, title_color=NAVY,
         sub_color=RGBColor.from_string("5B4A1A"))
textbox(s, Inches(5.85), Inches(1.5), Inches(7.0), Inches(5.6), bullets([
    ("130,502 parameters — a deliberately compact architecture suited to commodity hardware.", {"bold": True}),
    ("The element-wise product F_A ⊙ F_B is the key inductive bias: at every time position the "
     "network directly sees “are these two moving together or opposite?” — correlation evidence "
     "without hand-coding correlation.", {}),
    ("Shared (siamese) encoder: both series are embedded by identical weights, so the comparison "
     "is symmetric by construction.", {}),
    ("Lag is a 21-class classification, not regression: true posteriors can be multimodal "
     "(periodic signals alias), and argmax + the full distribution give an estimate plus confidence.", {}),
    ("β needs no head: with z-scored inputs the OLS slope equals ρ, so β̂ = ρ̂·σB/σA at readout.", {}),
    ("Ablation-friendly: encoder depth, fusion terms and heads are all config-driven.", {}),
], size=13.5))

# ============================================================ 9. training
s = new_slide()
kicker_title(s, "Method", "Training Objective and Procedure")
textbox(s, Inches(0.62), Inches(1.45), Inches(5.3), Inches(5.7), bullets([
    ("L = CE_masked(lag) + 4 · MSE(ρ) — the lag term is masked for independent pairs "
     "(their lag is undefined).", {}),
    ("Data generated on the fly: 8,000 steps × batch 128 ≈ 1.02M unique pairs, no dataset on disk.", {}),
    ("AdamW (lr 2e-3, weight decay 1e-4) with cosine decay to zero.", {}),
    ("Model selection on a fixed 4,096-pair validation set (seed 123) by lag-accuracy-±1 minus ρ-MAE; "
     "best checkpoint at step 7,000. Test seed (999) never touched during training.", {}),
    ("Inference cost: 2.4 ms per pair (0.12 ms/pair batched) — a single forward pass replaces "
     "the per-pair lag scan.", {"bold": True}),
    ("Validation accuracy rises from the 4.8% random-guess floor to 72.5% exact / 82.9% within ±1.", {}),
]))
picture(s, FIG / "fig1_training_curves.png", Inches(6.1), Inches(1.7), Inches(7.0), Inches(5.2))

# ============================================================ 10. synthetic benchmark
s = new_slide()
kicker_title(s, "Results 1/4", "Synthetic Benchmark vs. Classical Baseline")
picture(s, FIG / "fig5_model_vs_baseline.png", Inches(0.55), Inches(1.6), Inches(7.6), Inches(5.3))
textbox(s, Inches(8.35), Inches(1.6), Inches(4.5), Inches(5.6), bullets([
    ("Held-out synthetic test set: 20,000 pairs, seed 999.", {}),
    ("LeadLagNet: 71.6% exact lag, 82.7% within ±1, ρ MAE 0.165.", {}),
    ("Classical scan: 94.3% / 95.1% / 0.042 — superior on this benchmark.", {"bold": True}),
    ("This is expected: cross-correlation is near-optimal for exactly the linear, stationary "
     "relationships the generator produces.", {}),
    ("On strong signals (|ρ| ≥ 0.5, 69% of pairs) the model reaches 91.9% within ±1.", {}),
    ("The remaining gap is analyzed on the next slide.", {"color": GRAY}),
], size=13.5))

# ============================================================ 11. error analysis
s = new_slide()
kicker_title(s, "Results 2/4", "Error Analysis")
picture(s, FIG / "fig6_lag_confusion.png", Inches(0.55), Inches(1.55), Inches(5.9), Inches(4.9))
picture(s, FIG / "fig7_difficulty_buckets.png", Inches(6.7), Inches(1.55), Inches(6.1), Inches(4.9))
textbox(s, Inches(0.62), Inches(6.6), Inches(12.1), Inches(0.7),
        [("Left: a clean diagonal — when the model errs, it is almost always by ±1 step.   "
          "Right: both methods collapse as |ρ| → 0; with 128 samples a weak lag is "
          "statistically unidentifiable for any estimator — a property of the task, not a defect of the model.",
          {"size": 13.5, "color": INK})])

# ============================================================ 12. real market
s = new_slide()
kicker_title(s, "Results 3/4", "Real Markets — 4 / 4 Known Structures Recovered")
picture(s, FIG / "fig8_case_studies.png", Inches(0.5), Inches(1.6), Inches(6.9), Inches(5.4))
picture(s, FIG / "fig9_pair_heatmap.png", Inches(7.5), Inches(1.6), Inches(5.4), Inches(4.6))
textbox(s, Inches(7.55), Inches(6.25), Inches(5.3), Inches(1.0),
        [("Trained on zero real data, the model recovers: S&P→TWII and S&P→N225 at τ=+1, "
          "VIX↔S&P at τ=0 (negative), BTC↔ETH at τ=0 — and maps all 91 pairs in 11 ms.",
          {"size": 12.5, "color": INK})])

# ============================================================ 13. OOD
s = new_slide()
kicker_title(s, "Results 4/4", "Out-of-Distribution Generalization")
cards = [
    ("TEST A — SEMI-REAL", "90.3%", "lag within ±1",
     "Real market returns as base signal with constructed labels. Fat tails and volatility "
     "clustering are absent from training, yet accuracy improves over the synthetic test."),
    ("TEST B — OOD SYNTHETIC", "89.8%", "lag within ±1",
     "GARCH(1,1) bases + Student-t(3) noise: unseen signal families, almost no degradation."),
    ("TEST C — UNSEEN MARKETS", "19 / 20", "checks within ±1",
     "KOSPI, ASX 200, SOL, EURUSD — tickers absent from the whole pipeline. "
     "S&P→KOSPI τ=+1 ✓, BTC↔SOL τ=0 ✓; one miss: S&P→ASX off by one day."),
]
x = Inches(0.62)
for title, big, unit, body in cards:
    card(s, x, Inches(1.55), Inches(3.95), Inches(3.6))
    textbox(s, x + Inches(0.25), Inches(1.75), Inches(3.45), Inches(3.2),
            [(title, {"size": 12, "bold": True, "color": AMBER}),
             (big, {"size": 34, "bold": True, "color": NAVY}),
             (unit, {"size": 11.5, "color": GRAY, "space_after": 8}),
             (body, {"size": 12})])
    x += Inches(4.15)
c = card(s, Inches(0.62), Inches(5.45), Inches(12.1), Inches(1.6), fill=NAVY, line_color=NAVY)
textbox(s, Inches(0.95), Inches(5.65), Inches(11.5), Inches(1.3),
        [("BONUS FINDING", {"size": 11.5, "bold": True, "color": AMBER}),
         ("Both the model and the classical method report DXY↔EURUSD at τ=+1 instead of the expected 0: "
          "the model surfaced a real timestamp offset between FX and index daily bars in the data source "
          "— a data property we ourselves had not noticed.", {"size": 13.5, "color": WHITE})])

# ============================================================ 14. discussion
s = new_slide()
kicker_title(s, "Discussion", "Why Learn What We Can Compute?")
cols = [
    ("THE CLASSICAL METHOD WINS THE LINEAR CASE", NAVY, WHITE,
     ["94.3% vs 71.6% exact lag on our benchmark, ρ MAE 4× lower.",
      "Comparable speed at this scale (0.07 ms/pair vectorized).",
      "The gap concentrates in the weak-correlation regime (slide 11)."]),
    ("WHAT THE MODEL ALREADY ADDS", WHITE, INK,
     ["A full posterior over lags — it can say “I am not sure”; a scan cannot.",
      "A validated methodology: synthetic supervision transfers to real markets "
       "with zero labels (4/4 case studies, 90% OOD).",
      "A differentiable module that can be embedded in larger end-to-end systems."]),
    ("WHERE ONLY LEARNING WORKS", AMBER, NAVY,
     ["Volatility spillovers, threshold effects, asymmetric responses: Pearson "
       "correlation is mathematically blind to them.",
      "Each would need a hand-crafted statistic — or the same LeadLagNet with "
       "nonlinear samples added to the generator.",
      "Validate on a problem with a known answer first; then extend to problems without one."]),
]
x = Inches(0.62)
for title, fill, txt, items in cols:
    card(s, x, Inches(1.55), Inches(3.95), Inches(5.0), fill=fill,
         line_color=LINE if fill == WHITE else fill)
    sub = RGBColor.from_string("C7D2E5") if fill == NAVY else (GRAY if fill == WHITE else RGBColor.from_string("5B4A1A"))
    textbox(s, x + Inches(0.25), Inches(1.8), Inches(3.45), Inches(4.5),
            [(title, {"size": 13, "bold": True, "color": txt, "space_after": 10})] +
            bullets([(i, {"color": txt}) for i in items], size=12.5))
    x += Inches(4.15)
textbox(s, Inches(0.62), Inches(6.75), Inches(12.1), Inches(0.6),
        [("Positioning: this project does not replace cross-correlation — it validates the only "
          "approach that still works where cross-correlation fails.",
          {"size": 14, "bold": True, "color": NAVY})])

# ============================================================ 15. demo + limitations
s = new_slide()
kicker_title(s, "System", "Live Demo, Limitations & Future Work")
card(s, Inches(0.62), Inches(1.5), Inches(6.2), Inches(2.5))
textbox(s, Inches(0.92), Inches(1.65), Inches(5.7), Inches(2.2),
        [("INTERACTIVE DEMO (FastAPI + single-file web UI)", {"size": 12, "bold": True, "color": AMBER}),
         ("•  Pick any two assets and a window end date → verdict in plain words, "
          "lag distribution vs classical curve, 2.4 ms per inference", {"size": 12.5}),
         ("•  Validation page: live synthetic quiz against known answers, with session score "
          "and the classical method answering alongside", {"size": 12.5}),
         ("•  English / Chinese UI; offline parquet snapshot (no live calls)", {"size": 12.5})])
picture(s, Path("assets/screenshot.png"), Inches(0.62), Inches(4.2), Inches(6.2), Inches(2.95))
textbox(s, Inches(7.05), Inches(1.5), Inches(5.8), Inches(3.0),
        [("LIMITATIONS", {"size": 12, "bold": True, "color": AMBER})] + bullets([
            ("ρ̂ saturates near ±1 on very strong pairs (tanh over-confidence) and "
             "under-estimates weak ones — needs calibration.", {}),
            ("Weak-signal lag accuracy still trails the classical scan (49% vs 87% within ±1 "
             "in the 0.2–0.5 band).", {}),
            ("One OOD miss: S&P→ASX judged τ=0 instead of +1.", {}),
        ], size=12.5))
textbox(s, Inches(7.05), Inches(4.25), Inches(5.8), Inches(2.9),
        [("FUTURE WORK", {"size": 12, "bold": True, "color": AMBER})] + bullets([
            ("Nonlinear generator families (volatility spillover, thresholds) — where the "
             "classical method cannot follow.", {}),
            ("Intraday crypto (1h/4h): same model, only the validation data changes.", {}),
            ("Time-varying lag tracking — crisis windows (2020-03) show the structure shifting.", {}),
            ("Output calibration; multi-asset attention (full CMAN) as the long-term goal.", {}),
        ], size=12.5))

# ============================================================ 16. conclusion
s = new_slide(dark=True)
textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4),
        [("CONCLUSION", {"size": 13, "bold": True, "color": AMBER})])
textbox(s, Inches(0.9), Inches(0.95), Inches(11.8), Inches(0.7),
        [("Summary of findings", {"size": 26, "bold": True, "color": WHITE})])
claims = [
    ("1", "Synthetic supervision transfers to real markets",
     "Trained exclusively on generated pairs, the model recovers known market structure: "
     "4/4 lead-lag case studies, ~90% out-of-distribution accuracy, and it surfaced a genuine "
     "timestamp artifact in the data source."),
    ("2", "A compact estimator is sufficient",
     "130,502 parameters; a single 2.4 ms forward pass replaces the per-pair lag scan and "
     "additionally yields a full posterior over lags."),
    ("3", "The framework extends beyond classical statistics",
     "Cross-correlation remains superior in the linear regime. The validated synthetic-supervision "
     "pipeline, however, applies unchanged to nonlinear relationship families where no closed-form "
     "statistic exists."),
]
y = Inches(2.0)
for num, head, body in claims:
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.5), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = AMBER; chip.line.fill.background(); chip.shadow.inherit = False
    tf = chip.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = NAVY; r.font.name = FONT
    textbox(s, Inches(1.7), y - Inches(0.05), Inches(11.0), Inches(1.15),
            [(head, {"size": 16, "bold": True, "color": WHITE}),
             (body, {"size": 12.5, "color": RGBColor.from_string("C9D3E6")})])
    y += Inches(1.35)
textbox(s, Inches(0.95), Inches(6.45), Inches(11.7), Inches(0.5),
        [("Code: data/ · synth/ · models/ · train.py · eval.py · eval_ood.py · figures.py · serve.py + web UI",
          {"size": 11, "color": RGBColor.from_string("9FB0CC")})])

# ============================================================ 17. Q&A appendix
s = new_slide()
kicker_title(s, "Appendix", "Anticipated Questions")
QA = [
    ("If the classical scan is more accurate, why learn at all?",
     "Linear, stationary relationships are its optimal regime. The learned estimator adds a full "
     "posterior over lags (uncertainty), and extends to nonlinear families by changing only the "
     "generator — no new statistic has to be derived."),
    ("Is “no annotation” accurate when training is supervised?",
     "Yes. All supervision comes from the generative process itself; no human labels and no "
     "real-market labels are used at any stage. This is standard synthetic supervision "
     "(sim-to-real transfer)."),
    ("How is look-ahead bias avoided?",
     "Splits are applied to the window END date; normalization is per-window only; the test seed "
     "and the 2026 evaluation periods were never used for training or model selection."),
    ("The proposal mentioned attention — why a CNN?",
     "For a pairwise, fixed-length task a compact CNN is sufficient and easier to validate. "
     "Attention (the full CMAN architecture) is the planned extension for the multi-asset setting."),
    ("Why classify the lag instead of regressing it?",
     "The posterior can be multimodal — periodic signals alias across lags — and regression would "
     "average distinct peaks into a meaningless value. The softmax also provides uncertainty."),
    ("Can this be used for trading?",
     "Not directly. The model estimates relationship structure; turning a known lag into a "
     "profitable strategy involves transaction costs and market-efficiency questions outside "
     "this project's scope."),
]
qx, qy = Inches(0.62), Inches(1.5)
qw, qh = Inches(3.95), Inches(2.75)
for i, (q, a) in enumerate(QA):
    col, row = i % 3, i // 3
    x = qx + col * Inches(4.15)
    y = qy + row * Inches(2.95)
    card(s, x, y, qw, qh)
    textbox(s, x + Inches(0.22), y + Inches(0.18), qw - Inches(0.44), qh - Inches(0.36),
            [("Q: " + q, {"size": 11.5, "bold": True, "color": NAVY, "space_after": 6}),
             (a, {"size": 10.5, "color": INK})])

# ============================================================ speaker notes (zh-Hant)
NOTES_FULL = [  # 15-minute full script, kept for reference
    # 1
    "各位老師、同學好，我是徐崇恆。我的題目是 Cross-Market Lead-Lag Estimation for Financial Assets "
    "Using Deep Learning。這個專案處理的問題是：給定兩條金融時間序列，在沒有人工標註的情況下，"
    "用神經網路估計它們之間的 lead-lag（領先-落後）關係，包括誰領先誰、領先幾天、相關的強度和方向。"
    "接下來我會先介紹背景，再說明方法和模型設計，最後是驗證結果。",
    # 2
    "先說明背景。市場之間常有 lead-lag 關係：一個市場先反應，另一個市場之後跟著反應。傳統的做法是 "
    "cross-correlation（交叉相關）掃描：把序列 A 平移 τ 步，逐一計算和 B 的 correlation，取絕對值最大的"
    "位置。這個方法的限制是它假設關係為線性、平穩，而且每換一種假設就要重新設計統計量。右邊這張圖是"
    "真實資料：S&P 500 對三個亞洲指數的 cross-correlation，在 τ 等於正一的位置有明顯的峰值，向前的 "
    "correlation 是 0.42，反方向只有 0.004。原因是時區：美股收盤之後亞股才開盤，資訊的傳遞被交易時間"
    "隔開。這個專案想確認的是：神經網路能不能直接從兩段原始序列學會估計這種關係。",
    # 3
    "問題定義。輸入是兩條對齊的窗口，各 128 個交易日的日報酬，各自做 z-score 正規化，所以模型看到的是"
    "形狀，不是尺度。輸出有三項：第一是 lag，設計成 21 類的分類，範圍正負 10 步；第二是 correlation ρ，"
    "介於正負一之間，正負號代表方向，絕對值低於 0.2 視為不相關；第三是 scale β，因為輸入已經 z-score，"
    "迴歸係數在數學上等於 ρ，所以 β 不需要額外的輸出頭，由 ρ 和兩個窗口的標準差換算即可。lag 用分類"
    "而不用回歸，原因是後驗分佈可能是多峰的，例如週期性訊號會產生混疊；分類的 softmax 也能同時表達"
    "不確定性。",
    # 4
    "方法的核心是 synthetic supervision（合成監督）。實際的訓練集是用合成的方式生成的：先產生序列 A，"
    "把 A 平移 τ 步、乘上 β、加上強度 σ 的雜訊，得到 B。因為 B 是這樣構造出來的，τ、ρ、β 在生成的當下"
    "就已經知道。所以訓練在形式上是監督式的，但不需要任何人工標註，也沒有用到真實市場的標籤。"
    "這樣做的原因有三個：第一，標註成本為零；第二，資料可以一直生成、不需要儲存，模型不會過擬合某個"
    "固定的資料集；第三，難度可以控制，σ 決定 correlation 的強弱，樣本能覆蓋從很明顯到幾乎看不出來的"
    "整個範圍。訓練完成後，用真實市場資料驗證，這些資料在訓練過程完全沒有出現過。",
    # 5
    "系統分成六個階段：資料下載、資料驗證、生成器、模型、評估、demo。所有階段共用同一份 config，"
    "包含資產清單、窗口長度、lag 範圍和切分日期。時間切分以窗口的結束日為準：訓練資料只到 2025 年底，"
    "2026 年第一季是 validation，第二季是 test，不做跨時間的洗牌，正規化也不使用未來的資訊。另外一個"
    "決定是 cross-correlation baseline 在做模型之前就先完成，它同時是對照組，也是後續評估的參考答案"
    "來源。",
    # 6
    "資料是 14 個資產、十年的日線，涵蓋美股、VIX、亞洲與歐洲股指、加密貨幣、商品、利率和美元指數。"
    "這裡有一個對齊上的問題：加密貨幣一週交易七天，如果直接取所有日期的聯集當日曆，股市序列在週末會"
    "被填入大量為零的報酬，影響 correlation 的估計。處理方式是把日曆限制在週一到週五，加密貨幣週末的"
    "變動併入週一。假日的缺值最多 forward-fill 三天，更長的缺口維持缺值、相關窗口直接跳過。資料層共有 "
    "86 項自動化檢查，包括在建模之前，先用已知的市場關係確認對齊是正確的。",
    # 7
    "生成器的構造公式如投影片所示。基底序列有三種：白噪音、隨機係數的 AR(1)、和多個頻率的正弦疊加，"
    "用意是避免模型只學到單一類型的形狀。τ 在正負十之間均勻抽樣；β 的絕對值在 0.3 到 3 之間、正負各半；"
    "σ 控制 correlation 的強弱。另外保留 15% 完全獨立的序列對，讓模型學會判斷不相關的情況，這些樣本的 "
    "lag 沒有定義，對應的損失會被遮罩。訓練之前先做了一個一致性檢查：用傳統方法還原生成器的標籤，"
    "lag 的精確率是 94%，表示標籤和評估的約定是一致的。",
    # 8
    "模型是 siamese（孿生）架構的 1D CNN，參數量約十三萬。兩條序列經過同一組權重的 encoder，"
    "比較因此是對稱的。中間的 correlation fusion 是設計上的重點：把兩組特徵 F_A、F_B，和它們的逐元素"
    "乘積一起拼接。加入乘積項的原因是，網路在每個時間位置可以直接看到兩條序列是同向還是反向，"
    "相當於把 correlation 的證據提供給網路，而不需要手工去計算。後面接兩層卷積、pooling 和全連接層，"
    "最後是兩個輸出頭：lag 的 21 類 softmax，和 ρ 的 tanh。",
    # 9
    "損失函數是 lag 的 cross-entropy 加上 ρ 的 MSE，權重一比四。獨立序列對的 lag 沒有定義，"
    "它們的 cross-entropy 項會被遮罩。訓練資料即時生成，總量約一百萬對。optimizer 用 AdamW，"
    "學習率餘弦衰減；模型選擇用一個固定的 validation set，依 lag 準確率和 ρ 誤差的綜合分數挑選。"
    "最終 validation 的 lag 精確率是 72.5%，容差正負一是 82.9%，相較之下隨機猜測是 4.8%。"
    "推論只需要一次 forward pass，不需要逐 lag 掃描。",
    # 10
    "結果的第一部分是合成 test set，兩萬筆，沒有用於訓練或選模。模型的數字是 71.6%、82.7%、ρ 誤差 "
    "0.165；傳統方法是 94.3%、95.1%、0.042。在這個 benchmark 上傳統方法比較好，這點如實報告。"
    "這個結果在預期之內：cross-correlation 對線性、平穩的關係本來就接近最優，而生成器產生的正是"
    "這種關係。在訊號較強的子集，也就是 |ρ| 大於 0.5 的約七成樣本上，模型可以到 91.9%。"
    "差距的來源在下一頁分析。",
    # 11
    "左邊是 lag 的 confusion matrix，分佈集中在對角線，表示模型即使錯，大多也只差一步。右邊把樣本按 "
    "correlation 強度分組：強訊號區模型 92%、baseline 是 99.8%；中等訊號區是 49% 對 87%；|ρ| 低於 0.2 的"
    "弱訊號區兩者都低，17% 和 23%。這說明在 128 個樣本點之下，弱相關的 lag 對任何估計方法都接近"
    "不可辨識，是任務本身的限制。模型與 baseline 的差距，集中在中等與弱訊號這兩區。",
    # 12
    "真實市場的驗證。模型沒有用任何真實資料訓練，測試用四組已知關係：S&P 500 對台股和日經，已知領先"
    "一天；VIX 對 S&P 500，已知同日反相關；BTC 對 ETH，已知同日高度正相關。四組的 lag 和方向都與"
    "已知結果以及 baseline 一致。右邊的 heatmap 是模型對全部 91 個資產對的 correlation 估計，"
    "可以看到股指之間的正相關群、VIX 整列的負相關、和加密貨幣的群聚，整體結構合理。",
    # 13
    "再來是 out-of-distribution（分佈外）的測試，三項的內容都不在訓練分佈內。Test A 用真實市場報酬"
    "當基底，由我們施加已知的 τ 和 β。真實資料的肥尾和波動聚集是訓練時沒有的，結果準確率 90.3%，"
    "比合成 test set 還高一些，原因是真實日報酬接近白噪音，lag 的歧義比較少。Test B 用 GARCH 基底加 "
    "Student-t 雜訊，89.8%，幾乎沒有退化。Test C 用完全沒出現過的市場：KOSPI、ASX、SOL、EURUSD，"
    "20 個檢查點對了 19 個。另外有一個觀察：模型和 baseline 都回報 DXY 對 EURUSD 的 lag 是正一，"
    "而不是預期的零。查證之後，這是資料源裡外匯和指數的日線收盤時間戳不同步造成的，"
    "模型反映的是資料中實際存在的偏移。",
    # 14
    "這一頁回答一個需要回答的問題：傳統方法在線性情況下更準，為什麼還要用學習的方法。左欄是事實："
    "在我們的 benchmark 上傳統方法較好，速度也相當。中欄是模型目前的價值：它輸出整個 lag 的機率分佈，"
    "可以表達不確定，這是點估計做不到的；另外 synthetic supervision 到真實市場的遷移已經得到驗證。"
    "右欄是用學習方法的主要理由：波動率傳染、門檻效應這類非線性關係，Pearson correlation 在數學上"
    "測不到；傳統做法是為每種假設設計新的統計量，而這個框架只需要在生成器裡加入對應的樣本。"
    "我們的定位是：先在有已知答案的問題上驗證方法，再延伸到沒有封閉解的問題。",
    # 15
    "系統最後整合成一個互動的 demo：FastAPI 後端加網頁前端，選兩個資產和窗口結束日，會顯示模型的判斷、"
    "lag 的機率分佈和 baseline 的對照；另外有一頁可以現場生成已知答案的合成題目，驗證模型的行為。"
    "侷限有三點：ρ 在極強相關時會飽和，需要 calibration；弱訊號區的 lag 準確率落後 baseline；"
    "OOD 測試有一個未命中。未來工作：非線性的生成器、加密貨幣的盤中頻率、時變關係的追蹤，"
    "以及 attention 架構，也就是 proposal 裡 CMAN 的方向。",
    # 16
    "總結三點。第一，synthetic supervision 可以遷移到真實市場：純合成資料訓練的模型，在四組已知的 "
    "lead-lag 關係上與事實一致，OOD 準確率約九成，也找出了資料源的一個時間戳偏移。第二，小模型已經"
    "足夠：約十三萬參數，一次 forward pass 取代逐 lag 掃描，並且附帶不確定性分佈。第三，這個框架的"
    "延伸性：cross-correlation 在線性情況下仍然較強，但對它測不到的非線性關係，這套已驗證的流程可以"
    "直接套用。我的報告到這裡，謝謝大家。",
    # 17
    "這一頁是備用的問答整理，不一定會放出來，視提問情況使用。最可能被問的是第一題：傳統方法更準，"
    "為什麼還要用學習的方法。回答的重點是：線性情況本來就是傳統方法的最優範圍；模型的價值在於"
    "輸出不確定性，以及只要改生成器就能延伸到傳統統計量測不到的非線性關係。第二常見的是標註的問題："
    "訓練確實是監督式的，但所有監督訊號來自生成過程，沒有人工標註。其餘問題依卡片內容回答即可。",
]
# 6-minute condensed script (the active speaker notes)
NOTES = [
    # 1
    "各位老師、同學好，我是徐崇恆，題目是 Cross-Market Lead-Lag Estimation for Financial Assets Using "
    "Deep Learning。問題是：給定兩條金融時間序列，在沒有人工標註的情況下，用神經網路估計它們的 "
    "lead-lag 關係——誰領先誰、領先幾天、相關多強。",
    # 2
    "背景。市場之間常有 lead-lag 關係，傳統做法是 cross-correlation 掃描：逐一平移、逐一算 correlation、"
    "取最大值。限制是假設線性、平穩，而且每換一種假設就要重新設計統計量。右圖是真實資料：S&P 500 領先"
    "亞洲指數一個交易日，原因是時區。我們想確認的是，神經網路能不能直接從序列學會估計這種關係。",
    # 3
    "輸入是兩條 128 天、各自 z-score 的報酬窗口。輸出三項：lag，21 類分類；correlation ρ；scale β，"
    "由 ρ 換算、不需要額外的頭。lag 用分類不用回歸，原因是後驗可能多峰，softmax 也能表達不確定性。",
    # 4
    "方法的核心是 synthetic supervision。訓練集是合成的：把 A 平移 τ 步、乘上 β、加雜訊得到 B，"
    "所以標籤在生成當下就知道。訓練是監督式的，但沒有任何人工標註。好處是標註成本為零、資料可以"
    "一直生成、難度可控。真實資料只用來驗證，訓練時完全沒出現過。",
    # 5
    "系統六個階段，共用一份 config。時間切分以窗口結束日為準：訓練到 2025 年底，2026 年 Q1 是 "
    "validation、Q2 是 test，不跨時間洗牌。baseline 在建模之前先完成，作為對照組和參考答案。",
    # 6
    "資料是 14 個資產、十年日線。對齊上有一個問題：加密貨幣一週交易七天，直接取聯集會讓股市序列在"
    "週末多出大量零報酬，所以日曆限制在週一到週五。資料層共 86 項自動化檢查。",
    # 7
    "生成器的基底有三種：白噪音、AR(1)、正弦疊加，避免模型只學到單一形狀；15% 的樣本完全獨立，"
    "讓模型學會回答不相關。訓練前用傳統方法驗證過標籤的一致性。",
    # 8
    "模型是 siamese 的 1D CNN，約十三萬參數。設計重點是 correlation fusion：把兩組特徵和它們的"
    "逐元素乘積拼接，乘積項讓網路在每個時間位置直接看到兩條序列是同向還是反向。輸出兩個頭："
    "lag 的 softmax 和 ρ 的 tanh。",
    # 9
    "損失是 lag 的 cross-entropy 加 ρ 的 MSE，獨立對的 lag 損失會遮罩。validation 的 lag 精確率 "
    "72.5%、容差正負一 82.9%，隨機猜測是 4.8%。推論只需要一次 forward pass，不用逐 lag 掃描。",
    # 10
    "合成 test set 上，模型是 71.6% 和 82.7%；傳統方法 94.3% 和 95.1%，比較好，這點如實報告。"
    "這在預期內：線性、平穩正是它的最優範圍。在訊號較強的子集，模型可以到 91.9%。",
    # 11
    "誤差分析兩個重點：confusion matrix 顯示模型即使錯，大多只差一步；按強度分組後，弱訊號區"
    "兩種方法都低，因為在 128 個樣本點下，弱相關的 lag 本來就接近不可辨識，是任務本身的限制。",
    # 12
    "真實市場驗證。四組已知關係：S&P 500 對台股和日經領先一天、VIX 同日反相關、BTC 對 ETH 同日"
    "正相關——模型全部與已知結果一致，而它沒有看過任何真實資料。右邊的 heatmap 是 91 個資產對的"
    "整體結構，股指群、VIX 負相關、加密貨幣群都合理。",
    # 13
    "Out-of-distribution 測試三項：真實報酬當基底 90.3%；GARCH 加肥尾雜訊 89.8%；完全沒見過的市場，"
    "20 個檢查點對 19 個。另外模型和 baseline 都回報 DXY 對 EURUSD 有一天的偏移，查證後是資料源"
    "時間戳不同步，模型反映的是資料的實際狀況。",
    # 14
    "傳統方法更準，為什麼還要用學習的方法？線性本來就是它的最優範圍。模型的價值有二：輸出整個 lag "
    "的機率分佈，能表達不確定；以及非線性關係 Pearson correlation 測不到，傳統做法要重新設計統計量，"
    "這個框架只需要改生成器。先在有答案的問題上驗證，再延伸到沒有封閉解的問題。",
    # 15
    "系統整合成一個互動 demo，可以即時比較模型和 baseline。侷限：ρ 在極強相關時會飽和、弱訊號區"
    "落後 baseline。未來工作：非線性生成器、盤中頻率、attention 架構。",
    # 16
    "總結三點：synthetic supervision 可以遷移到真實市場；小模型已經足夠，一次 forward pass 並附帶"
    "不確定性分佈；框架可以延伸到傳統統計測不到的關係。報告到這裡，謝謝大家。",
    # 17
    "備用問答頁，視提問使用。最可能的問題是第一題：傳統方法更準為什麼還要學習——回答重點是線性是"
    "它的最優範圍，模型的價值在不確定性和非線性的延伸性。",
]

assert len(NOTES_FULL) == slide_no, f"{len(NOTES_FULL)} full notes vs {slide_no} slides"
assert len(NOTES) == slide_no, f"{len(NOTES)} notes vs {slide_no} slides"
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

OUT.parent.mkdir(exist_ok=True)
prs.save(OUT)
print(f"saved {OUT} ({slide_no} slides, notes on all)")

script_path = OUT.parent / "speech_script.txt"
parts = []
for i, note in enumerate(NOTES, 1):
    parts.append(f"--- Slide {i} ---")
    parts.append(note)
    parts.append("")
script_path.write_text("\n".join(parts), encoding="utf-8")
print(f"saved {script_path}")
